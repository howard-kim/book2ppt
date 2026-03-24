"""
Structured data → PowerPoint generator

Clones the 5 template slides for each unit and updates the content text boxes.
Template slide structure (by index):
  0 – Passage  : content in "TextBox 2"
  1 – Q1       : content in "TextBox 2"
  2 – Q2       : content in "TextBox 2"
  3 – Mission  : content in "TextBox 2"
  4 – Syntax   : content in "TextBox 4"
"""

import copy
from pathlib import Path
from typing import Dict, List

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

_FIELDS = ["passage", "q1", "q2", "q3", "mission", "syntax"]
_CONTENT_BOX = ["TextBox 2", "TextBox 2", "TextBox 2", "TextBox 2", "TextBox 2", "TextBox 4"]
_TEMPLATE_INDEX = [0, 1, 2, 2, 3, 4]
_CHAPTER_SLIDE_BOX = "TextBox 2"


def generate_ppt(
    data: List[Dict[str, str]],
    template_path: str,
    output_path: str,
) -> None:
    tp = Path(template_path)
    if not tp.exists():
        raise FileNotFoundError(f"Template not found: {tp}")

    prs = Presentation(str(tp))

    if len(prs.slides) < 5:
        raise ValueError(
            f"template.pptx must have at least 5 slides, got {len(prs.slides)}"
        )

    # Keep references to the 5 template slides before we add new ones
    template_slides = [prs.slides[i] for i in range(5)]
    chapter_counts: Dict[str, int] = {}
    chapter_numbers: Dict[str, int] = {}

    for absolute_no, unit in enumerate(data, start=1):
        chapter = unit.get("chapter", "default")
        if chapter not in chapter_numbers:
            chapter_numbers[chapter] = len(chapter_numbers) + 1
            _add_chapter_slide(prs, template_slides[0], chapter_numbers[chapter], chapter)
        chapter_counts[chapter] = chapter_counts.get(chapter, 0) + 1
        unit_no = chapter_counts[chapter]

        print(
            f"[INFO] Generating slides for unit {absolute_no} / {len(data)} "
            f"(chapter={chapter}, number={unit_no}) …"
        )
        badge_applied = False
        for slide_idx, (field, box_name) in enumerate(zip(_FIELDS, _CONTENT_BOX)):
            content = unit.get(field, "").strip()
            if not content:
                continue
            template_idx = _TEMPLATE_INDEX[slide_idx]
            new_slide = _clone_slide(prs, template_slides[template_idx])
            _set_text(new_slide, box_name, content)
            if not badge_applied:
                _set_existing_unit_badge(new_slide, unit_no)
                badge_applied = True

    # Remove the original 5 template slides (now at index 0-4)
    for _ in range(5):
        _delete_slide(prs, 0)

    prs.save(str(output_path))
    print(f"[INFO] Saved → {output_path}  ({len(prs.slides)} slides total)")


def _clone_slide(prs: Presentation, source) -> object:
    """Append a deep copy of source slide to the presentation."""
    new_slide = prs.slides.add_slide(source.slide_layout)
    sp_tree = new_slide.shapes._spTree
    for elem in list(sp_tree):
        sp_tree.remove(elem)
    rel_map = _clone_slide_relationships(source.part, new_slide.part)
    for elem in source.shapes._spTree:
        cloned = copy.deepcopy(elem)
        _rewrite_rel_ids(cloned, rel_map)
        sp_tree.append(cloned)
    return new_slide


def _clone_slide_relationships(source_part, target_part) -> dict[str, str]:
    """
    Clone non-layout relationships so copied shapes keep working.

    Slide XML can reference images and hyperlinks by rId. If we deep-copy only the
    shape XML, those references point to relationships that don't exist on the new
    slide, which breaks embedded SVG/image assets.
    """
    rel_map: dict[str, str] = {}
    allowed_reltypes = {
        RT.IMAGE,
        RT.HYPERLINK,
    }

    for rel in source_part.rels.values():
        if rel.reltype not in allowed_reltypes:
            continue

        if rel.is_external:
            new_rid = target_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rid = target_part.relate_to(rel.target_part, rel.reltype)

        rel_map[rel.rId] = new_rid

    return rel_map


def _rewrite_rel_ids(shape_xml, rel_map: dict[str, str]) -> None:
    """Rewrite relationship ids in copied shape XML to match the new slide."""
    rel_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    for elem in shape_xml.iter():
        for attr_name, attr_value in list(elem.attrib.items()):
            if not attr_name.startswith(f"{{{rel_ns}}}"):
                continue
            if attr_value in rel_map:
                elem.set(attr_name, rel_map[attr_value])


def _set_text(slide, box_name: str, text: str) -> None:
    """Replace all text in a named text box."""
    for shape in _iter_shapes(slide.shapes):
        if shape.name == box_name and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            lines = text.split("\n") if text else [""]
            for line_no, line in enumerate(lines):
                para = tf.paragraphs[0] if line_no == 0 else tf.add_paragraph()
                para.add_run().text = line
            return
    print(f"[WARN] Text box '{box_name}' not found on slide")


def _set_existing_unit_badge(slide, unit_no: int) -> None:
    badge_text = str(unit_no)
    sp_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

    for shape_el in slide._element.findall(f".//{{{sp_ns}}}sp"):
        cNvPr = shape_el.find(f"./{{{sp_ns}}}nvSpPr/{{{sp_ns}}}cNvPr")
        if cNvPr is None:
            continue

        shape_name = cNvPr.get("name", "")
        text_nodes = shape_el.findall(f".//{{{a_ns}}}t")
        shape_text = "".join(node.text or "" for node in text_nodes).strip()

        if shape_name == "TextBox 9" or shape_text == "#":
            if text_nodes:
                text_nodes[0].text = badge_text
                for extra in text_nodes[1:]:
                    extra.text = ""
                return

            body = shape_el.find(f".//{{{sp_ns}}}txBody")
            if body is not None:
                para = body.find(f"./{{{a_ns}}}p")
                if para is not None:
                    run = para.find(f"./{{{a_ns}}}r")
                    if run is not None:
                        text_node = run.find(f"./{{{a_ns}}}t")
                        if text_node is not None:
                            text_node.text = badge_text
                            return

    for shape in _iter_shapes(slide.shapes):
        if _is_badge_text_shape(shape):
            shape.text = badge_text
            return

    print("[WARN] Unit badge text box not found on passage slide")


def _add_chapter_slide(prs: Presentation, source_slide, chapter_no: int, chapter_name: str) -> None:
    chapter_slide = _clone_slide(prs, source_slide)
    _remove_named_shapes(chapter_slide, {"Group 7"})
    _set_text(chapter_slide, _CHAPTER_SLIDE_BOX, f"Chapter {chapter_no}\n{chapter_name}")


def _is_badge_text_shape(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False

    text = shape.text.strip()
    if shape.name == "TextBox 9":
        return True
    if text in {"#", "1"} and shape.left < 200000 and shape.top < 100000:
        return True
    if shape.width < 900000 and shape.height < 900000 and shape.left < 200000:
        return text in {"#", "1", ""}
    return False


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)


def _remove_named_shapes(slide, shape_names: set[str]) -> None:
    for shape in list(slide.shapes):
        if shape.name in shape_names:
            slide.shapes._spTree.remove(shape._element)


def _delete_slide(prs: Presentation, slide_index: int) -> None:
    """Remove a slide by index from the presentation."""
    xml_slides = prs.slides._sldIdLst
    sldId = xml_slides[slide_index]
    rId = sldId.get(qn("r:id"))
    xml_slides.remove(sldId)
    prs.part.drop_rel(rId)
